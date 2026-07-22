import logging
from io import BytesIO

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".txt", ".docx", ".pdf", ".pptx"}


def extract_text(buf: BytesIO, filename: str) -> str:
    """Fayldan matn chiqaradi. Bo'sh satr qaytarsa — o'qib bo'lmadi."""
    ext = ("." + filename.rsplit(".", 1)[-1]).lower() if "." in filename else ""

    if ext == ".txt":
        return _read_txt(buf)
    if ext == ".docx":
        return _read_docx(buf)
    if ext == ".pdf":
        return _read_pdf(buf)
    if ext == ".pptx":
        return _read_pptx(buf)

    raise ValueError(f"Qo'llab-quvvatlanmaydigan format: {ext or 'nomaʼlum'}")


def _read_txt(buf: BytesIO) -> str:
    raw = buf.read()
    for encoding in ("utf-8-sig", "utf-8", "cp1251"):
        try:
            return raw.decode(encoding).strip()
        except UnicodeDecodeError:
            continue
    raise ValueError("Fayl kodlashini o'qib bo'lmadi (UTF-8 yoki CP1251 kerak).")


def _read_docx(buf: BytesIO) -> str:
    try:
        from docx import Document  # python-docx
        doc = Document(buf)
        lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t:
                        lines.append(t)
        return "\n".join(lines)
    except Exception:
        logger.exception("_read_docx xato")
        raise ValueError("DOCX faylni o'qib bo'lmadi.")


def _read_pdf(buf: BytesIO) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(buf)
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(pages).strip()
    except Exception:
        logger.exception("_read_pdf xato")
        raise ValueError("PDF faylni o'qib bo'lmadi.")


def _read_pptx(buf: BytesIO) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(buf)
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        lines.append(t)
        return "\n".join(lines)
    except Exception:
        logger.exception("_read_pptx xato")
        raise ValueError("PPTX faylni o'qib bo'lmadi.")
