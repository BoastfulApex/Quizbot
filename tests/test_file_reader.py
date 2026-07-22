from io import BytesIO

import pytest

from bot.utils.file_reader import extract_text


# --- TXT ---

def test_extract_text_txt_utf8():
    buf = BytesIO("Salom dunyo!".encode("utf-8"))
    assert extract_text(buf, "test.txt") == "Salom dunyo!"


def test_extract_text_txt_cp1251():
    buf = BytesIO("Привет мир".encode("cp1251"))
    assert "Привет" in extract_text(buf, "test.txt")


def test_extract_text_txt_utf8_bom():
    buf = BytesIO("BOM fayl".encode("utf-8-sig"))
    assert extract_text(buf, "test.TXT") == "BOM fayl"


# --- DOCX ---

def test_extract_text_docx():
    from docx import Document
    doc = Document()
    doc.add_paragraph("Birinchi qator")
    doc.add_paragraph("Ikkinchi qator")
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)

    result = extract_text(buf, "mavzu.docx")
    assert "Birinchi qator" in result
    assert "Ikkinchi qator" in result


def test_extract_text_docx_with_table():
    from docx import Document
    doc = Document()
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Katak A"
    table.rows[0].cells[1].text = "Katak B"
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)

    result = extract_text(buf, "table.docx")
    assert "Katak A" in result
    assert "Katak B" in result


# --- PDF ---

def test_extract_text_pdf():
    from pypdf import PdfWriter
    writer = PdfWriter()
    page = writer.add_blank_page(width=200, height=200)
    buf = BytesIO()
    writer.write(buf)
    buf.seek(0)

    result = extract_text(buf, "empty.pdf")
    assert isinstance(result, str)


# --- PPTX ---

def test_extract_text_pptx():
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "Slayd matni"
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)

    result = extract_text(buf, "prezentatsiya.pptx")
    assert "Slayd matni" in result


# --- Noto'g'ri kengaytma ---

def test_extract_text_unsupported_extension_raises():
    buf = BytesIO(b"data")
    with pytest.raises(ValueError, match="Qo'llab-quvvatlanmaydigan"):
        extract_text(buf, "fayl.xls")
